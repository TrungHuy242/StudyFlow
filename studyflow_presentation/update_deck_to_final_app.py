from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parent
SRC = ROOT / "StudyFlow_Professional_Figma_Deck.pptx"
OUT = ROOT / "StudyFlow_Professional_Final_App_Deck.pptx"


REPLACEMENTS = {
    "Ứng dụng hỗ trợ sinh viên quản lý lịch học, deadline và kế hoạch ôn tập một cách chủ động.":
        "Ứng dụng đã hoàn thiện, hỗ trợ sinh viên quản lý lịch học, deadline và kế hoạch ôn tập một cách chủ động.",
    "Hiện trạng: các module StudyFlow đã có UI và logic chính; có thể đóng gói, tích hợp dữ liệu thật và phát hành.":
        "Hoàn thành: StudyFlow đã có UI, logic chính và các luồng demo mobile; sẵn sàng trình bày, bàn giao và mở rộng.",
    "Các bước tiếp theo sau khi app StudyFlow đã hoàn thiện chức năng cốt lõi.":
        "Mở rộng sau hoàn thiện",
    "Professional deck with Figma-based mobile visuals":
        "Final app presentation • Figma design + Flutter implementation",
    "StudyFlow gom lịch học, deadline và kế hoạch ôn tập vào một trải nghiệm mobile thống nhất.":
        "StudyFlow là ứng dụng mobile đã hoàn thiện, gom lịch học, deadline và kế hoạch ôn tập vào một trải nghiệm thống nhất.",
    "Tạo công cụ giúp sinh viên học có kế hoạch, đúng hạn và ít căng thẳng hơn.":
        "Xây dựng ứng dụng hoàn chỉnh giúp sinh viên học có kế hoạch, đúng hạn và ít căng thẳng hơn.",
    "Những nỗi đau chính được ưu tiên giải quyết trong phiên bản đầu.":
        "Những nỗi đau chính đã được StudyFlow giải quyết trong sản phẩm hoàn thiện.",
    "Ý tưởng được tạo từ pain points và chuyển thành feature ưu tiên.":
        "Ý tưởng được chuyển hóa thành các chức năng đã triển khai trong StudyFlow.",
    "Prototype kiểm tra luồng tác vụ cốt lõi trước khi đầu tư high-fi.":
        "Low-Fi prototype dùng để kiểm chứng luồng cốt lõi trước khi hoàn thiện giao diện.",
    "Kiểm tra khả năng hiểu flow và thao tác nhanh trên mobile.":
        "Kiểm tra và xác nhận các luồng chính dễ hiểu, thao tác nhanh trên mobile.",
    "Thay đổi sau vòng test để tăng rõ ràng và tốc độ thao tác.":
        "Các cải tiến đã được áp dụng vào bản app hoàn thiện để tăng rõ ràng và tốc độ thao tác.",
    "Prototype hoàn thiện mô phỏng đầy đủ các trạng thái người dùng sẽ gặp.":
        "High-Fi prototype thể hiện đầy đủ các trạng thái và được chuyển thành app Flutter hoàn chỉnh.",
    "Kịch bản demo đề xuất cho phần trình bày.":
        "Kịch bản demo app hoàn thiện trong phần trình bày.",
    "Prototype Flutter đã được tách thành các app/folder riêng để dễ phát triển tiếp.":
        "Ứng dụng Flutter đã được triển khai với cấu trúc module rõ ràng, dễ bảo trì và mở rộng.",
    "Hướng phát triển: gộp các module vào một StudyFlow app duy nhất, thêm lưu trữ local/backend và notification.":
        "Hiện trạng: các module StudyFlow đã có UI và logic chính; có thể đóng gói, tích hợp dữ liệu thật và phát hành.",
    "Phiên bản hiện tại chứng minh được trải nghiệm cốt lõi của StudyFlow.":
        "Sản phẩm hoàn thiện đã đáp ứng đầy đủ trải nghiệm cốt lõi của StudyFlow.",
    "2 app Flutter độc lập cho Schedule và Deadline từ Figma.":
        "Các module Flutter cho Schedule, Deadline và luồng học tập đã được hoàn thiện từ Figma.",
    "App flutter hoàn chình":
        "App Flutter hoàn chỉnh",
    "Bài học chính trong quá trình chuyển Figma thành prototype Flutter.":
        "Bài học chính trong quá trình chuyển Figma thành app Flutter hoàn thiện.",
    "Các bước tiếp theo để biến prototype thành sản phẩm StudyFlow hoàn chỉnh.":
        "Các bước tiếp theo sau khi app StudyFlow đã hoàn thiện chức năng cốt lõi.",
    "Gộp Schedule + Deadline + Review Plan vào một codebase Flutter.":
        "Tinh chỉnh một codebase Flutter thống nhất cho Schedule, Deadline và Review Plan.",
    "Thêm lưu trữ local bằng SQLite/Hive và đồng bộ cloud khi đăng nhập.":
        "Bổ sung dữ liệu người dùng thật, lưu trữ local và đồng bộ cloud khi đăng nhập.",
    "Tích hợp notification cho lớp học, deadline và phiên ôn tập.":
        "Hoàn thiện notification cho lớp học, deadline và phiên ôn tập.",
    "Import lịch từ Google Calendar/LMS hoặc file syllabus.":
        "Tích hợp import lịch từ Google Calendar/LMS hoặc file syllabus.",
    "Gợi ý kế hoạch ôn tập tự động theo deadline, độ khó và thời gian trống.":
        "Nâng cấp gợi ý kế hoạch ôn tập tự động theo deadline, độ khó và thời gian trống.",
    "Chạy usability test vòng 2 với sinh viên thật và đo time-on-task.":
        "Chạy usability test sau phát hành với sinh viên thật và đo time-on-task.",
    "Review":
        "Ôn tập",
}


TABLE_REPLACEMENTS = {
    "Ngày/tuần/tháng, thêm/sửa lịch học, chi tiết lớp":
        "Đã hoàn thiện xem ngày/tuần/tháng, thêm/sửa lịch học, chi tiết lớp",
    "Deadline list, today, week, overdue, add/edit/detail":
        "Đã hoàn thiện danh sách, hôm nay, tuần, quá hạn, thêm/sửa/chi tiết",
    "Định nghĩa dữ liệu ScheduleItem, DeadlineItem, priority, state":
        "Định nghĩa dữ liệu, trạng thái và logic cho lịch học/deadline",
    "UI, navigation, form validation, interaction":
        "UI hoàn chỉnh, điều hướng, validation form và tương tác người dùng",
    "Đảm bảo project compile và test pass":
        "Đã kiểm tra compile, analyze và test pass",
}


def replace_in_text_frame(text_frame, replacements):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            for _ in range(2):
                for old, new in replacements.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)


def update_table(table):
    for row in table.rows:
        for cell in row.cells:
            replace_in_text_frame(cell.text_frame, REPLACEMENTS)
            replace_in_text_frame(cell.text_frame, TABLE_REPLACEMENTS)


def main():
    prs = Presentation(SRC)
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                replace_in_text_frame(shape.text_frame, REPLACEMENTS)
            if getattr(shape, "has_table", False):
                update_table(shape.table)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
