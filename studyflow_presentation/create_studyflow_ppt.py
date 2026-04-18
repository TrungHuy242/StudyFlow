from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("StudyFlow_Tong_quan_du_an.pptx")
W, H = Inches(13.333), Inches(7.5)

PRIMARY = RGBColor(52, 120, 246)
GREEN = RGBColor(34, 197, 94)
ORANGE = RGBColor(245, 158, 11)
RED = RGBColor(239, 68, 68)
PINK = RGBColor(236, 72, 153)
CYAN = RGBColor(6, 182, 212)
INK = RGBColor(16, 24, 40)
MUTED = RGBColor(102, 112, 133)
LIGHT = RGBColor(247, 248, 252)
CARD = RGBColor(255, 255, 255)
BORDER = RGBColor(229, 231, 239)
BLUE_SOFT = RGBColor(234, 242, 255)
GREEN_SOFT = RGBColor(235, 255, 243)
ORANGE_SOFT = RGBColor(255, 246, 221)
RED_SOFT = RGBColor(255, 241, 242)


def bg(slide, color=LIGHT):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def text(slide, x, y, w, h, value, size=16, color=INK, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = value
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def bullets(slide, x, y, w, h, items, size=15, color=MUTED, bullet=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}" if bullet else item
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def shape(slide, x, y, w, h, fill=CARD, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(kind, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = BORDER
    s.line.width = Pt(0.75)
    return s


def pill(slide, x, y, w, h, value, fill=BLUE_SOFT, color=PRIMARY, size=10):
    s = shape(slide, x, y, w, h, fill)
    s.line.color.rgb = fill
    tf = s.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return s


def title(slide, heading, subtitle="", eyebrow="", page=0):
    if eyebrow:
        text(slide, Inches(0.65), Inches(0.42), Inches(4.0), Inches(0.25),
             eyebrow.upper(), 9, PRIMARY, True)
    text(slide, Inches(0.65), Inches(0.74), Inches(8.9), Inches(0.55),
         heading, 26, INK, True)
    if subtitle:
        text(slide, Inches(0.65), Inches(1.28), Inches(9.8), Inches(0.34),
             subtitle, 12, MUTED)
    if page:
        text(slide, Inches(0.55), Inches(7.05), Inches(2.0), Inches(0.2),
             "StudyFlow", 8, MUTED, True)
        text(slide, Inches(11.75), Inches(7.05), Inches(0.8), Inches(0.2),
             f"{page:02d}", 8, MUTED, True, PP_ALIGN.RIGHT)


def icon(slide, x, y, value, fill=PRIMARY):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.42), Inches(0.42))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = fill
    tf = s.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    r.font.name = "Arial"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)
    return s


def phone(slide, x, y, w, h, heading, rows, accent=PRIMARY):
    shape(slide, x, y, w, h, CARD)
    text(slide, x + Inches(0.18), y + Inches(0.18), w - Inches(0.36),
         Inches(0.25), heading, 12, INK, True)
    yy = y + Inches(0.62)
    for i, row in enumerate(rows):
        fill = BLUE_SOFT if i == 0 else RGBColor(248, 250, 252)
        shape(slide, x + Inches(0.18), yy, w - Inches(0.36), Inches(0.42), fill)
        text(slide, x + Inches(0.32), yy + Inches(0.12), w - Inches(0.64),
             Inches(0.14), row, 7.5, accent if i == 0 else MUTED, i == 0)
        yy += Inches(0.5)


def table(slide, x, y, w, h, rows, widths):
    tbl = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h).table
    for i, width in enumerate(widths):
        tbl.columns[i].width = width
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            cell.margin_left = Pt(4)
            cell.margin_right = Pt(4)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY if r == 0 else CARD
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Arial"
            p.font.size = Pt(8.4 if r else 9)
            p.font.bold = r == 0
            p.font.color.rgb = RGBColor(255, 255, 255) if r == 0 else INK
    return tbl


def cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, RGBColor(242, 247, 255))
    text(s, Inches(0.72), Inches(0.82), Inches(5.0), Inches(0.28),
         "UX/UI CASE STUDY", 10, PRIMARY, True)
    text(s, Inches(0.72), Inches(1.35), Inches(6.4), Inches(1.1),
         "StudyFlow", 50, INK, True)
    text(s, Inches(0.78), Inches(2.42), Inches(6.55), Inches(0.85),
         "Ứng dụng hỗ trợ sinh viên quản lý lịch học, deadline và kế hoạch ôn tập.",
         20, MUTED)
    pill(s, Inches(0.78), Inches(3.48), Inches(1.55), Inches(0.38), "Lịch học")
    pill(s, Inches(2.55), Inches(3.48), Inches(1.55), Inches(0.38), "Deadline", RED_SOFT, RED)
    pill(s, Inches(4.32), Inches(3.48), Inches(1.65), Inches(0.38), "Ôn tập", GREEN_SOFT, GREEN)
    phone(s, Inches(7.75), Inches(0.72), Inches(2.0), Inches(4.6), "Lịch học",
          ["UX/UI 07:00", "Web Dev 09:45", "Database 13:00"], PRIMARY)
    phone(s, Inches(10.0), Inches(1.2), Inches(2.0), Inches(4.6), "Deadline",
          ["Quiz hôm nay", "UX Research 2 ngày", "Project DB 5 ngày"], RED)
    text(s, Inches(0.78), Inches(6.62), Inches(4.7), Inches(0.24),
         "Prepared for StudyFlow project presentation", 10, MUTED)


def bullet_slide(prs, page, heading, subtitle, items, accent=PRIMARY):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, heading, subtitle, f"{page:02d}", page + 1)
    bullets(s, Inches(0.85), Inches(1.75), Inches(7.0), Inches(4.4),
            items, 17, INK, True)
    phone(s, Inches(8.75), Inches(1.45), Inches(2.4), Inches(4.85),
          heading[:18], ["Tổng quan", "Theo dõi", "Hoàn thành"], accent)


def cards_slide(prs, page, heading, subtitle, cards):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, heading, subtitle, f"{page:02d}", page + 1)
    for i, (h, b, fill, color) in enumerate(cards):
        x = Inches(0.75 + (i % 3) * 4.1)
        y = Inches(1.72 + (i // 3) * 1.72)
        shape(s, x, y, Inches(3.55), Inches(1.22), fill)
        icon(s, x + Inches(0.22), y + Inches(0.22), str(i + 1), color)
        text(s, x + Inches(0.78), y + Inches(0.22), Inches(2.55), Inches(0.24),
             h, 14, INK, True)
        bullets(s, x + Inches(0.78), y + Inches(0.55), Inches(2.5), Inches(0.5),
                [b], 10.8, MUTED, False)


def overview(prs):
    cards_slide(prs, 1, "Tổng quan dự án",
                "StudyFlow gom các việc học quan trọng vào một trải nghiệm mobile thống nhất.",
                [
                    ("Lịch học", "Xem ngày/tuần/tháng, thêm/sửa lớp, phòng học, giảng viên.", BLUE_SOFT, PRIMARY),
                    ("Deadline", "Theo dõi bài tập, mức ưu tiên, quá hạn, tiến độ hoàn thành.", RED_SOFT, RED),
                    ("Kế hoạch ôn tập", "Chia nhỏ mục tiêu học, nhắc ôn tập, theo dõi tiến độ.", GREEN_SOFT, GREEN),
                ])


def research(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "Nghiên cứu người dùng",
          "Tổng hợp insight định tính cho nhóm người dùng sinh viên đại học.", "04", 5)
    rows = [
        ["Phương pháp", "Mục đích", "Insight chính"],
        ["Interview 1:1", "Hiểu thói quen lập kế hoạch", "Sinh viên dùng nhiều công cụ nhưng bỏ sót cập nhật."],
        ["Survey nhanh", "Xác định tần suất quên deadline", "Deadline và lịch học là nhu cầu thường xuyên nhất."],
        ["Desk research", "So sánh app hiện có", "Ít app gom lịch học + deadline + ôn tập trong UX mobile gọn."],
        ["Usability test", "Kiểm tra flow thêm deadline/lịch", "CTA rõ và trạng thái quá hạn giúp người dùng hiểu nhanh."],
    ]
    table(s, Inches(0.7), Inches(1.65), Inches(11.95), Inches(3.35),
          rows, [Inches(2.2), Inches(3.15), Inches(6.6)])
    bullets(s, Inches(0.9), Inches(5.42), Inches(11.2), Inches(0.55),
            ["Ghi chú: nếu nhóm đã khảo sát thực tế, có thể thay phần này bằng số liệu thật."],
            12, MUTED, False)


def competitor(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "Phân tích đối thủ",
          "Ít nhất 20 app được chia theo nhóm cạnh tranh trực tiếp và gián tiếp.", "05", 6)
    rows = [
        ["Nhóm", "Ứng dụng", "Điểm mạnh", "Khoảng trống cho StudyFlow"],
        ["Student planner", "MyStudyLife, myHomework, Power Planner, iStudiez Pro, School Planner", "Theo dõi lớp, bài tập, nhắc nhở", "UX có thể cũ, chưa nhấn mạnh kế hoạch ôn tập cá nhân"],
        ["Deadline tracker", "Assignment Planner, Egenda, Class Timetable, Studious, Homework Planner", "Tập trung bài tập/hạn nộp", "Ít kết nối với lịch học và tiến độ ôn thi"],
        ["Productivity", "Todoist, TickTick, Microsoft To Do, Any.do, Trello", "Task mạnh, nhắc nhở, label", "Không sinh viên-hóa ngữ cảnh môn học/lớp học"],
        ["Calendar/planning", "Google Calendar, Notion, Structured, FlowSavvy, Apple Calendar", "Lịch, time blocking, template", "Cần tự cấu hình, dễ quá tải với sinh viên mới"],
        ["Study tools", "Quizlet, Anki, Forest, Focus To-Do, Google Classroom, Canvas Student", "Ôn tập, tập trung, LMS", "Không gom toàn bộ journey học tập cá nhân"],
    ]
    table(s, Inches(0.48), Inches(1.42), Inches(12.35), Inches(4.55),
          rows, [Inches(1.45), Inches(3.55), Inches(3.0), Inches(4.35)])
    text(s, Inches(0.75), Inches(6.15), Inches(11.5), Inches(0.34),
         "Cơ hội: StudyFlow nên là “student command center”, không chỉ là to-do list.",
         18, INK, True, PP_ALIGN.CENTER)


def persona(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "Persona", "Đại diện nhóm sinh viên cần quản lý nhiều môn và nhiều hạn nộp cùng lúc.", "07", 8)
    shape(s, Inches(0.78), Inches(1.55), Inches(3.0), Inches(4.7), BLUE_SOFT)
    text(s, Inches(1.05), Inches(1.9), Inches(2.4), Inches(0.42), "Minh Anh", 24, INK, True, PP_ALIGN.CENTER)
    text(s, Inches(1.05), Inches(2.42), Inches(2.4), Inches(0.3), "Sinh viên năm 2 ngành CNTT", 12, MUTED, False, PP_ALIGN.CENTER)
    bullets(s, Inches(1.05), Inches(3.05), Inches(2.35), Inches(1.4),
            ["Nộp bài đúng hạn", "Biết hôm nay cần học gì", "Không ôn thi sát giờ"], 12, INK, True)
    text(s, Inches(4.25), Inches(1.65), Inches(3.6), Inches(0.3), "Hành vi hiện tại", 16, INK, True)
    bullets(s, Inches(4.25), Inches(2.05), Inches(3.7), Inches(1.8),
            ["Ghi deadline trong chat nhóm hoặc Google Calendar.", "Dùng to-do list nhưng ít cập nhật tiến độ.", "Thường học theo cảm giác thay vì kế hoạch."], 13, MUTED, True)
    text(s, Inches(8.35), Inches(1.65), Inches(3.6), Inches(0.3), "Nhu cầu", 16, INK, True)
    bullets(s, Inches(8.35), Inches(2.05), Inches(3.8), Inches(1.8),
            ["Một màn hình biết việc nào gấp.", "Flow thêm deadline nhanh.", "Nhắc nhở nhẹ, không gây áp lực.", "Theo dõi tiến độ rõ bằng phần trăm."], 13, MUTED, True)
    shape(s, Inches(4.25), Inches(4.55), Inches(7.9), Inches(1.1), ORANGE_SOFT)
    text(s, Inches(4.55), Inches(4.83), Inches(7.2), Inches(0.42),
         "“Mình không thiếu app, mình thiếu một chỗ cho biết nên làm gì tiếp theo.”",
         17, INK, True, PP_ALIGN.CENTER)


def journey(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "User Journey", "Hành trình từ khi mở app đến khi hoàn thành một nhiệm vụ học tập.", "08", 9)
    steps = [
        ("Mở app", "Biết hôm nay có gì", "Lo lắng"),
        ("Xem lịch", "Kiểm tra lớp và giờ trống", "Ổn định"),
        ("Chọn deadline", "Biết việc gấp", "Tập trung"),
        ("Lập ôn tập", "Chia nhỏ bài học", "Chủ động"),
        ("Hoàn thành", "Cập nhật tiến độ", "Nhẹ nhõm"),
    ]
    for i, (h, b, mood) in enumerate(steps):
        x = Inches(0.65 + i * 2.48)
        shape(s, x, Inches(1.9), Inches(2.05), Inches(2.75), CARD)
        icon(s, x + Inches(0.18), Inches(2.12), str(i + 1), PRIMARY)
        text(s, x + Inches(0.18), Inches(2.72), Inches(1.7), Inches(0.3), h, 13, INK, True)
        bullets(s, x + Inches(0.18), Inches(3.12), Inches(1.75), Inches(0.58), [b], 10.5, MUTED, False)
        pill(s, x + Inches(0.18), Inches(4.05), Inches(1.1), Inches(0.28), mood, BLUE_SOFT, PRIMARY, 8.5)
    text(s, Inches(1.0), Inches(5.38), Inches(11.3), Inches(0.36),
         "Design implication: home screen phải trả lời nhanh “hôm nay cần làm gì?”",
         18, INK, True, PP_ALIGN.CENTER)


def low_fi(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "Low-Fi Design",
          "Wireframe tập trung vào cấu trúc nội dung và flow chính trước khi xử lý visual.", "10", 11)
    phone(s, Inches(0.9), Inches(1.55), Inches(2.55), Inches(4.85),
          "Lịch học", ["Segment: Ngày/Tuần/Tháng", "Class card", "Bottom nav"], PRIMARY)
    phone(s, Inches(5.35), Inches(1.55), Inches(2.55), Inches(4.85),
          "Deadline", ["Search + filter", "Overdue section", "Upcoming cards"], RED)
    phone(s, Inches(9.8), Inches(1.55), Inches(2.55), Inches(4.85),
          "Ôn tập", ["Goal", "Session list", "Progress"], GREEN)


def prototype(prs, page, heading, subtitle, flow_items):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, heading, subtitle, f"{page:02d}", page + 1)
    for i, item in enumerate(flow_items):
        x = Inches(0.9 + (i % 4) * 3.0)
        y = Inches(2.0 + (i // 4) * 1.35)
        shape(s, x, y, Inches(2.15), Inches(0.72), CARD)
        text(s, x, y + Inches(0.24), Inches(2.15), Inches(0.18),
             item, 11.5, INK, True, PP_ALIGN.CENTER)
        if i < len(flow_items) - 1 and i % 4 != 3:
            line = s.shapes.add_connector(1, x + Inches(2.15), y + Inches(0.36),
                                          x + Inches(2.95), y + Inches(0.36))
            line.line.color.rgb = PRIMARY
            line.line.width = Pt(1.5)


def testing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "Usability Testing", "Kiểm tra khả năng hiểu flow và thao tác nhanh trên mobile.", "12", 13)
    rows = [
        ["Task", "Tiêu chí thành công", "Kết quả/Insight"],
        ["Thêm deadline mới", "Hoàn thành dưới 60 giây", "Cần label rõ cho ngày/giờ đến hạn"],
        ["Tìm deadline hôm nay", "Không quá 2 lần tap", "Tab Hôm nay giúp giảm thời gian tìm kiếm"],
        ["Sửa tiến độ", "Người dùng hiểu % hoàn thành", "Progress bar nên đi kèm số %"],
        ["Xem quá hạn", "Nhận biết ngay rủi ro", "Màu đỏ và badge Quá hạn hoạt động tốt"],
    ]
    table(s, Inches(0.7), Inches(1.55), Inches(11.95), Inches(3.65),
          rows, [Inches(2.7), Inches(3.15), Inches(6.1)])
    bullets(s, Inches(0.9), Inches(5.55), Inches(10.9), Inches(0.45),
            ["Khuyến nghị: ưu tiên trạng thái rõ, CTA thêm mới nổi bật, giảm nhập liệu bằng dropdown và date/time picker."],
            14, INK, False)


def high_fi(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "High-Fi Design",
          "Visual direction dựa trên Figma: mobile-first, card rõ, màu trạng thái dễ nhận biết.", "14", 15)
    phone(s, Inches(0.75), Inches(1.35), Inches(2.25), Inches(4.9),
          "Schedule", ["Ngày/Tuần/Tháng", "UX/UI Design", "Database Systems"], PRIMARY)
    phone(s, Inches(3.35), Inches(1.35), Inches(2.25), Inches(4.9),
          "Deadlines", ["Quá hạn", "Assignment UX", "Quiz hôm nay"], RED)
    phone(s, Inches(5.95), Inches(1.35), Inches(2.25), Inches(4.9),
          "Detail", ["Progress 60%", "Chỉnh sửa", "Hoàn thành"], GREEN)
    text(s, Inches(8.8), Inches(1.65), Inches(3.3), Inches(0.35),
         "Design System", 18, INK, True)
    bullets(s, Inches(8.8), Inches(2.1), Inches(3.5), Inches(2.3),
            ["8px radius cho card/button", "Primary blue #3478F6", "Red cho quá hạn", "Green cho hoàn thành", "Typography rõ trên mobile"], 13, MUTED, True)
    pill(s, Inches(8.8), Inches(4.75), Inches(0.9), Inches(0.34), "#3478F6", BLUE_SOFT, PRIMARY, 8.5)
    pill(s, Inches(9.85), Inches(4.75), Inches(0.9), Inches(0.34), "#EF4444", RED_SOFT, RED, 8.5)
    pill(s, Inches(10.9), Inches(4.75), Inches(0.9), Inches(0.34), "#22C55E", GREEN_SOFT, GREEN, 8.5)


def code_impl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "Code Implementation",
          "Prototype Flutter đã được tách thành các app/folder riêng để dễ phát triển tiếp.", "18", 19)
    rows = [
        ["Module", "File/Folder", "Vai trò"],
        ["Schedule Calendar", "schedule_calendar_app", "Ngày/tuần/tháng, thêm/sửa lịch học, chi tiết lớp"],
        ["Assignments Deadlines", "assignments_deadlines_app", "Deadline list, today, week, overdue, add/edit/detail"],
        ["Model layer", "models/*.dart", "Định nghĩa dữ liệu ScheduleItem, DeadlineItem, priority, state"],
        ["Screen layer", "screens/*.dart", "UI, navigation, form validation, interaction"],
        ["Quality check", "flutter analyze/test/build web", "Đảm bảo project compile và test pass"],
    ]
    table(s, Inches(0.65), Inches(1.5), Inches(12.05), Inches(4.35),
          rows, [Inches(2.0), Inches(3.0), Inches(7.05)])
    text(s, Inches(0.85), Inches(6.1), Inches(11.4), Inches(0.34),
         "Hướng phát triển: gộp các module vào một StudyFlow app duy nhất, thêm lưu trữ local/backend và notification.",
         15, INK, True, PP_ALIGN.CENTER)


def qa(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    text(s, Inches(0.95), Inches(1.8), Inches(11.4), Inches(0.9),
         "Q&A", 56, INK, True, PP_ALIGN.CENTER)
    text(s, Inches(1.55), Inches(2.95), Inches(10.2), Inches(0.4),
         "Cảm ơn thầy/cô và các bạn đã lắng nghe.", 20, MUTED, False, PP_ALIGN.CENTER)
    pill(s, Inches(5.25), Inches(3.75), Inches(2.8), Inches(0.42),
         "StudyFlow", BLUE_SOFT, PRIMARY, 14)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    cover(prs)
    overview(prs)
    bullet_slide(prs, 2, "Vấn đề", "Sinh viên cần học chủ động hơn nhưng thông tin bị phân tán.",
                 ["Lịch học, bài tập, nhắc nhở và tài liệu nằm ở nhiều app khác nhau.",
                  "Deadline đến gần nhưng không được ưu tiên rõ theo ngày/tuần.",
                  "Ôn tập thường bị động, chỉ bắt đầu khi sát kiểm tra.",
                  "Quá nhiều việc nhỏ khiến người dùng không biết bắt đầu từ đâu."], RED)
    bullet_slide(prs, 3, "Mục tiêu", "Tạo công cụ giúp sinh viên học có kế hoạch, đúng hạn và ít căng thẳng hơn.",
                 ["Giảm việc quên lịch học, deadline và phiên ôn tập quan trọng.",
                  "Hiển thị ưu tiên rõ ràng theo hôm nay, tuần này và quá hạn.",
                  "Cho phép thêm/sửa/xóa nhanh trên mobile với ít bước thao tác.",
                  "Tạo cảm giác tiến bộ bằng trạng thái, phần trăm hoàn thành và checklist.",
                  "Thiết kế dễ tiếp cận: chữ rõ, tương phản tốt, vùng chạm đủ lớn."], GREEN)
    research(prs)
    competitor(prs)
    cards_slide(prs, 6, "Pain Points", "Những nỗi đau chính được ưu tiên giải quyết trong phiên bản đầu.",
                [("Quên deadline", "Không thấy mức độ khẩn cấp trước khi quá muộn.", RED_SOFT, RED),
                 ("Không biết ưu tiên", "Không rõ hôm nay nên học gì trước.", ORANGE_SOFT, ORANGE),
                 ("Khó cập nhật lịch", "Lịch học thay đổi nhưng thao tác sửa mất thời gian.", BLUE_SOFT, PRIMARY),
                 ("Bài lớn bị mơ hồ", "Không chia nhỏ tiến độ thành các bước làm được.", GREEN_SOFT, GREEN),
                 ("Ôn tập thiếu kế hoạch", "Thiếu nhắc nhở và mục tiêu theo ngày.", BLUE_SOFT, CYAN),
                 ("Giao diện quá tải", "Nhiều chữ làm người dùng bỏ qua thông tin chính.", RED_SOFT, PINK)])
    persona(prs)
    journey(prs)
    bullet_slide(prs, 9, "Ideation", "Ý tưởng được tạo từ pain points và chuyển thành feature ưu tiên.",
                 ["HMW giúp sinh viên thấy deadline quan trọng trước khi quá muộn?",
                  "HMW kết nối lịch học với kế hoạch làm bài và ôn tập?",
                  "HMW biến tiến độ học thành tín hiệu trực quan, dễ hiểu?",
                  "Feature ưu tiên: calendar view, deadline badges, progress bar, review plan."], PRIMARY)
    low_fi(prs)
    prototype(prs, 11, "Low-Fi Prototype", "Prototype kiểm tra luồng tác vụ cốt lõi trước khi đầu tư high-fi.",
              ["Home", "Lịch học", "Thêm lịch", "Deadline", "Chi tiết", "Ôn tập", "Kế hoạch", "Hoàn thành"])
    testing(prs)
    cards_slide(prs, 13, "Cải tiến", "Những thay đổi sau vòng test để tăng rõ ràng và tốc độ thao tác.",
                [("Trước", "Deadline chỉ là danh sách dài, khó biết việc nào gấp.", RED_SOFT, RED),
                 ("Sau", "Tách Quá hạn, Sắp tới, Hôm nay, Tuần.", GREEN_SOFT, GREEN),
                 ("Trước", "Form nhập nhiều text tự do.", RED_SOFT, RED),
                 ("Sau", "Dùng dropdown, date picker, time picker, chip ưu tiên.", GREEN_SOFT, GREEN)])
    high_fi(prs)
    prototype(prs, 15, "High-Fi Prototype", "Prototype hoàn thiện mô phỏng đầy đủ các trạng thái người dùng sẽ gặp.",
              ["Schedule flow", "Deadline flow", "Review flow", "Empty state", "Add/Edit", "Detail", "Progress", "Complete"])
    cards_slide(prs, 16, "Accessibility", "Thiết kế để sinh viên thao tác nhanh và đọc được trong nhiều bối cảnh.",
                [("Contrast", "Màu chữ chính đậm, trạng thái dùng màu + text badge.", BLUE_SOFT, PRIMARY),
                 ("Touch target", "Button, icon và chip ưu tiên khoảng 40-52px.", GREEN_SOFT, GREEN),
                 ("Readable copy", "Nhãn ngắn: Hôm nay, Quá hạn, Thêm Deadline.", ORANGE_SOFT, ORANGE),
                 ("Navigation", "Bottom nav nhất quán cho các module StudyFlow.", BLUE_SOFT, PRIMARY),
                 ("Error prevention", "Date/time picker, dropdown môn học, validation form.", GREEN_SOFT, GREEN),
                 ("State feedback", "Snack bar, badge, progress bar và empty state rõ.", RED_SOFT, RED)])
    bullet_slide(prs, 17, "Demo mobile", "Kịch bản demo đề xuất cho phần trình bày.",
                 ["Mở app và xem deadline chưa hoàn thành.",
                  "Chuyển tab Hôm nay để thấy Quiz HTML/CSS.",
                  "Mở chi tiết deadline và tăng tiến độ.",
                  "Thêm một deadline mới với ngày, giờ, mức ưu tiên.",
                  "Chuyển sang lịch học để xem ngày/tuần/tháng.",
                  "Kết luận bằng trạng thái quá hạn và kế hoạch ôn tập."], RED)
    code_impl(prs)
    cards_slide(prs, 19, "Kết quả cuối", "Phiên bản hiện tại chứng minh được trải nghiệm cốt lõi của StudyFlow.",
                [("Hoàn thành", "2 app Flutter độc lập cho Schedule và Deadline từ Figma.", GREEN_SOFT, GREEN),
                 ("Tương tác", "Thêm/sửa/xóa, lọc, tìm kiếm, cập nhật tiến độ, empty state.", BLUE_SOFT, PRIMARY),
                 ("Kiểm chứng", "Analyzer sạch, widget test pass, build web thành công.", GREEN_SOFT, GREEN),
                 ("Giá trị", "Dashboard học tập rõ hơn và ít bỏ sót việc quan trọng.", ORANGE_SOFT, ORANGE)])
    bullet_slide(prs, 20, "Reflection", "Bài học chính trong quá trình chuyển Figma thành prototype Flutter.",
                 ["Thiết kế tốt cần trạng thái dữ liệu rõ: empty, overdue, today, upcoming.",
                  "Sinh viên cần quyết định nhanh, nên ưu tiên hierarchy và badge trạng thái.",
                  "Figma giúp định hình visual, nhưng code cần tách model/screen để dễ mở rộng.",
                  "Nên validate với người dùng thật để thay giả định bằng số liệu đáng tin cậy."], GREEN)
    bullet_slide(prs, 21, "Next Steps", "Các bước tiếp theo để biến prototype thành sản phẩm StudyFlow hoàn chỉnh.",
                 ["Gộp Schedule + Deadline + Review Plan vào một codebase Flutter.",
                  "Thêm lưu trữ local bằng SQLite/Hive và đồng bộ cloud khi đăng nhập.",
                  "Tích hợp notification cho lớp học, deadline và phiên ôn tập.",
                  "Import lịch từ Google Calendar/LMS hoặc file syllabus.",
                  "Gợi ý kế hoạch ôn tập tự động theo deadline, độ khó và thời gian trống.",
                  "Chạy usability test vòng 2 với sinh viên thật và đo time-on-task."], PRIMARY)
    qa(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
