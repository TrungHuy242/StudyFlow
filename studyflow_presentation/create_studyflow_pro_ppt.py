from pathlib import Path
from textwrap import wrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "pro_assets"
OUT = ROOT / "StudyFlow_Professional_Figma_Deck.pptx"
ASSETS.mkdir(exist_ok=True)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
NAVY = RGBColor(12, 23, 43)
INK = RGBColor(16, 24, 40)
MUTED = RGBColor(102, 112, 133)
SOFT = RGBColor(246, 248, 252)
CARD = RGBColor(255, 255, 255)
BORDER = RGBColor(225, 231, 240)
BLUE = RGBColor(52, 120, 246)
GREEN = RGBColor(34, 197, 94)
RED = RGBColor(239, 68, 68)
ORANGE = RGBColor(245, 158, 11)
PURPLE = RGBColor(124, 92, 255)
PINK = RGBColor(236, 72, 153)
CYAN = RGBColor(6, 182, 212)


def font(size, bold=False):
    candidates = [
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/segoeuib.ttf" if bold else "C:/Windows/Fonts/segoeui.ttf",
    ]
    for path in candidates:
        try:
            return ImageFont.truetype(path, size)
        except Exception:
            continue
    return ImageFont.load_default()


def rgb(color):
    return (color[0], color[1], color[2])


def pil_color(c):
    return (c.rgb[0], c.rgb[1], c.rgb[2]) if hasattr(c, "rgb") else c


def draw_wrapped(draw, xy, text, width, fill, size=18, bold=False, spacing=4):
    x, y = xy
    f = font(size, bold)
    chars = max(8, int(width / (size * 0.55)))
    for line in wrap(text, chars):
        draw.text((x, y), line, fill=fill, font=f)
        y += size + spacing
    return y


def rr(draw, box, radius, fill, outline=None, width=1):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def phone_canvas(bg=(247, 248, 252)):
    im = Image.new("RGB", (375, 812), bg)
    d = ImageDraw.Draw(im)
    d.text((26, 15), "9:41", fill=(18, 24, 38), font=font(15, True))
    d.text((295, 15), "▮▮  ◔  ▰", fill=(18, 24, 38), font=font(11, True))
    return im, d


def bottom_nav(d, active="Việc"):
    rr(d, (2, 726, 373, 812), 0, (255, 255, 255), (230, 234, 242))
    items = [("Trang chủ", "⌂"), ("Lịch", "□"), ("Việc", "☑"), ("Thống kê", "▥"), ("Tôi", "○")]
    x_positions = [43, 119, 187, 260, 333]
    for (label, ic), x in zip(items, x_positions):
        color = (52, 120, 246) if label == active else (139, 149, 166)
        d.text((x - 8, 742), ic, fill=color, font=font(17, True))
        d.text((x - 22, 767), label, fill=color, font=font(10, True))
    rr(d, (118, 797, 257, 802), 3, (20, 24, 31))


def progress(d, x, y, w, pct, color):
    rr(d, (x, y, x + w, y + 5), 3, (231, 235, 244))
    rr(d, (x, y, x + int(w * pct), y + 5), 3, color)


def save_screen(name, im):
    path = ASSETS / f"{name}.png"
    im.save(path, optimize=True)
    return path


def deadline_card(d, y, title, subject, badge, pct, color, overdue=False):
    fill = (255, 241, 242) if overdue else (255, 255, 255)
    rr(d, (22, y, 353, y + 118), 12, fill, (254, 202, 202) if overdue else (229, 231, 239))
    rr(d, (39, y + 17, 79, y + 57), 10, tuple(int(v * 0.14 + 255 * 0.86) for v in color), None)
    d.text((50, y + 27), "▣", fill=color, font=font(16, True))
    d.text((90, y + 16), title, fill=(16, 24, 40), font=font(16, True))
    d.text((90, y + 40), subject, fill=(102, 112, 133), font=font(13, True))
    badge_color = (239, 68, 68) if overdue else ((245, 158, 11) if "Hôm nay" in badge else (34, 197, 94))
    rr(d, (270, y + 18, 335, y + 43), 10, tuple(int(v * 0.13 + 255 * 0.87) for v in badge_color), None)
    d.text((282, y + 24), badge, fill=badge_color, font=font(9, True))
    d.text((106, y + 75), "23:59", fill=(102, 112, 133), font=font(12, True))
    d.text((306, y + 75), f"{int(pct * 100)} %", fill=(102, 112, 133), font=font(12, True))
    progress(d, 92, y + 98, 244, pct, color)


def schedule_card(d, y, title, time, room, label, color):
    rr(d, (22, y, 353, y + 114), 12, (255, 255, 255), (229, 231, 239))
    rr(d, (39, y + 24, 45, y + 84), 3, color)
    d.text((61, y + 20), title, fill=(16, 24, 40), font=font(16, True))
    d.text((296, y + 20), time.split("-")[0].strip(), fill=(102, 112, 133), font=font(12, True))
    d.text((80, y + 52), time, fill=(102, 112, 133), font=font(12, True))
    d.text((205, y + 52), room, fill=(102, 112, 133), font=font(12, True))
    rr(d, (61, y + 82, 135, y + 104), 10, tuple(int(v * 0.14 + 255 * 0.86) for v in color), None)
    d.text((70, y + 87), label, fill=color, font=font(10, True))


def make_deadline_list():
    im, d = phone_canvas()
    d.text((21, 56), "Deadline", fill=(16, 24, 40), font=font(29, True))
    d.text((22, 96), "6 deadline chưa hoàn thành", fill=(102, 112, 133), font=font(14, True))
    rr(d, (312, 58, 352, 98), 12, (52, 120, 246))
    d.text((327, 68), "+", fill=(255, 255, 255), font=font(20, True))
    rr(d, (22, 128, 304, 168), 12, (255, 255, 255), (229, 231, 239))
    d.text((69, 141), "Tìm kiếm...", fill=(152, 162, 179), font=font(14))
    rr(d, (312, 128, 353, 168), 12, (255, 255, 255), (229, 231, 239))
    d.text((22, 204), "⚠ Quá hạn", fill=(239, 68, 68), font=font(15, True))
    deadline_card(d, 238, "Lab Report 3", "Database Systems", "Quá hạn", 0.0, (239, 68, 68), True)
    d.text((22, 386), "Sắp tới", fill=(16, 24, 40), font=font(17, True))
    deadline_card(d, 420, "Assignment 1 - UX", "UX/UI Design", "2 ngày", 0.6, (99, 102, 241))
    deadline_card(d, 548, "Quiz - HTML/CSS", "Web Development", "Hôm nay", 0.8, (34, 197, 94))
    bottom_nav(d, "Việc")
    return save_screen("figma_deadline_list", im)


def make_deadline_today():
    im, d = phone_canvas(bg=(255, 154, 11))
    d.text((22, 52), "Deadline hôm nay", fill=(255, 255, 255), font=font(22, True))
    d.text((22, 82), "Ngày 9 tháng 4, 2026", fill=(255, 245, 225), font=font(13, True))
    rr(d, (22, 128, 181, 204), 12, (255, 177, 55))
    rr(d, (193, 128, 352, 204), 12, (255, 177, 55))
    d.text((94, 142), "1", fill=(255, 255, 255), font=font(25, True))
    d.text((73, 172), "Cần làm", fill=(255, 255, 255), font=font(13, True))
    d.text((264, 142), "0", fill=(255, 255, 255), font=font(25, True))
    d.text((239, 172), "Khẩn cấp", fill=(255, 255, 255), font=font(13, True))
    rr(d, (22, 244, 353, 414), 14, (255, 255, 255))
    d.text((98, 260), "Quiz - HTML/CSS", fill=(16, 24, 40), font=font(17, True))
    d.text((99, 306), "Web Development", fill=(102, 112, 133), font=font(13, True))
    d.text((118, 338), "10:00", fill=(102, 112, 133), font=font(14, True))
    d.text((98, 372), "Tiến độ", fill=(16, 24, 40), font=font(14, True))
    d.text((307, 370), "80 %", fill=(102, 112, 133), font=font(12, True))
    progress(d, 100, 400, 237, .8, (52, 120, 246))
    bottom_nav(d, "Việc")
    return save_screen("figma_deadline_today", im)


def make_deadline_detail():
    im, d = phone_canvas()
    d.text((22, 56), "‹", fill=(102, 112, 133), font=font(26, True))
    rr(d, (22, 104, 353, 224), 14, (255, 255, 255), (229, 231, 239))
    rr(d, (38, 122, 78, 162), 10, (226, 232, 255))
    d.text((92, 121), "Assignment 1 - UX", fill=(16, 24, 40), font=font(17, True))
    d.text((92, 145), "UX/UI Design", fill=(102, 112, 133), font=font(13, True))
    d.text((306, 177), "60 %", fill=(102, 112, 133), font=font(12, True))
    progress(d, 92, 197, 244, .6, (52, 120, 246))
    rr(d, (22, 246, 181, 294), 12, (234, 242, 255), None)
    rr(d, (193, 246, 353, 294), 12, (235, 255, 243), None)
    d.text((54, 262), "Chưa hoàn thành", fill=(52, 120, 246), font=font(12, True))
    d.text((237, 262), "Hoàn tất", fill=(34, 197, 94), font=font(12, True))
    for y, label, value in [(324, "Ngày đến hạn", "2026-04-12"), (422, "Giờ đến hạn", "23:59"), (520, "Mô tả", "Hoàn thành bài nghiên cứu người dùng")]:
        rr(d, (22, y, 353, y + 76), 12, (255, 255, 255), (229, 231, 239))
        d.text((94, y + 16), label, fill=(102, 112, 133), font=font(13, True))
        d.text((94, y + 38), value, fill=(16, 24, 40), font=font(15, True))
    bottom_nav(d, "Việc")
    return save_screen("figma_deadline_detail", im)


def make_form_screen(name, heading, save_label, accent=BLUE):
    im, d = phone_canvas()
    d.text((22, 106), "‹", fill=(102, 112, 133), font=font(22, True))
    d.text((118, 110), heading, fill=(16, 24, 40), font=font(19, True))
    fields = [("Tiêu đề", "Assignment 1 - UX Research"), ("Môn học", "UX/UI Design"), ("Ngày đến hạn", "2026-04-12"), ("Giờ đến hạn", "23:59"), ("Mức độ ưu tiên", "Thấp     Bình thường     Khẩn cấp")]
    y = 180
    for label, value in fields:
        d.text((22, y), label, fill=(52, 64, 84), font=font(13, True))
        rr(d, (22, y + 28, 353, y + 82), 12, (255, 255, 255), (229, 231, 239))
        d.text((39, y + 45), value, fill=(102, 112, 133), font=font(13, True))
        y += 96
    rr(d, (2, 700, 373, 812), 0, (255, 255, 255), (229, 231, 239))
    rr(d, (22, 724, 353, 780), 12, pil_color(accent), None)
    d.text((132, 742), save_label, fill=(255, 255, 255), font=font(17, True))
    rr(d, (118, 797, 257, 802), 3, (20, 24, 31))
    return save_screen(name, im)


def make_deadline_overdue():
    im, d = phone_canvas(bg=(244, 45, 74))
    d.text((22, 52), "Quá hạn", fill=(255, 255, 255), font=font(25, True))
    d.text((22, 88), "1 deadline đã quá hạn", fill=(255, 227, 232), font=font(13, True))
    deadline_card(d, 160, "Lab Report 3", "Database Systems", "Quá hạn", 0, (239, 68, 68), False)
    rr(d, (22, 286, 353, 380), 14, (255, 255, 255))
    draw_wrapped(d, (39, 306), "Những deadline quá hạn sẽ ảnh hưởng đến điểm số và tiến độ học tập của bạn.", 285, (180, 35, 24), 13, True)
    bottom_nav(d, "Việc")
    return save_screen("figma_deadline_overdue", im)


def make_deadline_empty():
    im, d = phone_canvas()
    d.text((21, 100), "Deadline", fill=(16, 24, 40), font=font(28, True))
    rr(d, (147, 238, 227, 318), 40, (234, 242, 255))
    d.text((174, 260), "▣", fill=(52, 120, 246), font=font(26, True))
    d.text((111, 334), "Chưa có deadline", fill=(16, 24, 40), font=font(20, True))
    d.text((64, 370), "Thêm deadline để theo dõi các nhiệm vụ", fill=(102, 112, 133), font=font(13, True))
    d.text((106, 390), "và bài tập quan trọng", fill=(102, 112, 133), font=font(13, True))
    rr(d, (64, 434, 311, 482), 12, (52, 120, 246))
    d.text((107, 448), "Thêm deadline đầu tiên", fill=(255, 255, 255), font=font(14, True))
    bottom_nav(d, "Việc")
    return save_screen("figma_deadline_empty", im)


def make_schedule_day():
    im, d = phone_canvas()
    d.text((21, 108), "Lịch học", fill=(16, 24, 40), font=font(25, True))
    rr(d, (312, 104, 353, 145), 12, (52, 120, 246))
    d.text((327, 113), "+", fill=(255, 255, 255), font=font(20, True))
    rr(d, (22, 160, 353, 208), 12, (255, 255, 255), (229, 231, 239))
    rr(d, (27, 165, 128, 203), 10, (234, 242, 255))
    d.text((61, 174), "Ngày", fill=(52, 120, 246), font=font(13, True))
    d.text((171, 174), "Tuần", fill=(102, 112, 133), font=font(13, True))
    d.text((276, 174), "Tháng", fill=(102, 112, 133), font=font(13, True))
    rr(d, (2, 226, 373, 294), 0, (255, 255, 255), (229, 231, 239))
    d.text((164, 236), "Thứ 2", fill=(16, 24, 40), font=font(18, True))
    d.text((134, 260), "9 tháng 4, 2026", fill=(102, 112, 133), font=font(13, True))
    schedule_card(d, 310, "UX/UI Design", "07:00 - 09:30", "A301", "Lý thuyết", (124, 92, 255))
    schedule_card(d, 436, "Web Development", "09:45 - 12:15", "B202", "Lý thuyết", (34, 197, 94))
    schedule_card(d, 562, "Database Systems", "13:00 - 15:30", "C105", "Thực hành", (245, 158, 11))
    bottom_nav(d, "Lịch")
    return save_screen("figma_schedule_day", im)


def make_schedule_week():
    im, d = phone_canvas()
    d.text((21, 110), "Lịch tuần", fill=(16, 24, 40), font=font(24, True))
    rr(d, (22, 160, 353, 208), 12, (255, 255, 255), (229, 231, 239))
    rr(d, (137, 165, 239, 203), 10, (234, 242, 255))
    d.text((59, 174), "Ngày", fill=(102, 112, 133), font=font(13, True))
    d.text((170, 174), "Tuần", fill=(52, 120, 246), font=font(13, True))
    d.text((276, 174), "Tháng", fill=(102, 112, 133), font=font(13, True))
    d.text((129, 244), "Tuần 14 - 2026", fill=(16, 24, 40), font=font(18, True))
    days = ["T2\n6", "T3\n7", "T4\n8", "T5\n9", "T6\n10", "T7\n11", "CN\n12"]
    for i, day in enumerate(days):
        x = 48 + i * 46
        fill = (52, 120, 246) if "9" in day else (255, 255, 255)
        txt = (255, 255, 255) if "9" in day else (102, 112, 133)
        rr(d, (x, 284, x + 36, 330), 10, fill, (229, 231, 239))
        for j, part in enumerate(day.split("\n")):
            d.text((x + 10, 291 + j * 19), part, fill=txt, font=font(11, True))
    rr(d, (53, 342, 95, 397), 9, (232, 240, 255))
    d.text((57, 345), "UX/UI", fill=(52, 120, 246), font=font(10, True))
    rr(d, (100, 402, 160, 461), 9, (235, 255, 243))
    d.text((104, 404), "Web Dev", fill=(34, 197, 94), font=font(10, True))
    rr(d, (165, 530, 227, 589), 9, (255, 246, 221))
    d.text((168, 532), "Database", fill=(245, 158, 11), font=font(10, True))
    bottom_nav(d, "Lịch")
    return save_screen("figma_schedule_week", im)


def make_schedule_month():
    im, d = phone_canvas()
    d.text((21, 110), "Lịch tháng", fill=(16, 24, 40), font=font(24, True))
    d.text((129, 244), "Tháng 4, 2026", fill=(16, 24, 40), font=font(18, True))
    labels = ["T2", "T3", "T4", "T5", "T6", "T7", "CN"]
    for i, lab in enumerate(labels):
        d.text((38 + i * 45, 292), lab, fill=(102, 112, 133), font=font(11, True))
    day = 1
    for row in range(5):
        for col in range(7):
            x = 28 + col * 46
            y = 318 + row * 40
            if row == 0 and col < 2:
                continue
            if day > 30:
                continue
            active = day == 9
            rr(d, (x, y, x + 34, y + 34), 9, (52, 120, 246) if active else (255, 255, 255), (229, 231, 239))
            d.text((x + 10, y + 8), str(day), fill=(255, 255, 255) if active else (16, 24, 40), font=font(11, True))
            if day in [6, 7, 8, 9, 10, 11, 12]:
                rr(d, (x + 14, y + 27, x + 19, y + 32), 3, (34, 197, 94), None)
            day += 1
    d.text((22, 558), "Sự kiện ngày 9", fill=(16, 24, 40), font=font(17, True))
    schedule_card(d, 592, "UX/UI Design", "07:00 - 09:30", "A301", "Lý thuyết", (124, 92, 255))
    bottom_nav(d, "Lịch")
    return save_screen("figma_schedule_month", im)


def make_schedule_empty():
    im, d = phone_canvas()
    d.text((21, 102), "Lịch học", fill=(16, 24, 40), font=font(25, True))
    rr(d, (147, 238, 227, 318), 40, (234, 242, 255))
    d.text((173, 260), "□", fill=(52, 120, 246), font=font(28, True))
    d.text((114, 330), "Chưa có lịch học", fill=(16, 24, 40), font=font(20, True))
    d.text((64, 366), "Thêm lịch học để xem và quản lý thời", fill=(102, 112, 133), font=font(13, True))
    d.text((128, 386), "gian biểu của bạn", fill=(102, 112, 133), font=font(13, True))
    rr(d, (67, 430, 308, 478), 12, (52, 120, 246))
    d.text((110, 444), "Thêm lịch học đầu tiên", fill=(255, 255, 255), font=font(14, True))
    bottom_nav(d, "Lịch")
    return save_screen("figma_schedule_empty", im)


def make_lowfi(name, title_text, lines):
    im = Image.new("RGB", (375, 812), (250, 250, 250))
    d = ImageDraw.Draw(im)
    d.text((22, 42), title_text, fill=(80, 80, 80), font=font(23, True))
    rr(d, (22, 95, 353, 145), 8, (232, 232, 232), None)
    y = 180
    for i, line in enumerate(lines):
        rr(d, (22, y, 353, y + 92), 10, (255, 255, 255), (210, 210, 210))
        rr(d, (40, y + 18, 78, y + 56), 8, (225, 225, 225), None)
        d.text((96, y + 18), line, fill=(92, 92, 92), font=font(15, True))
        rr(d, (96, y + 50, 300, y + 58), 4, (225, 225, 225), None)
        y += 112
    bottom_nav(d, "Việc")
    return save_screen(name, im)


def generate_assets():
    paths = {
        "deadline_list": make_deadline_list(),
        "deadline_today": make_deadline_today(),
        "deadline_detail": make_deadline_detail(),
        "deadline_add": make_form_screen("figma_deadline_add", "Thêm Deadline", "Thêm Deadline", BLUE),
        "deadline_edit": make_form_screen("figma_deadline_edit", "Sửa Deadline", "Lưu", BLUE),
        "deadline_overdue": make_deadline_overdue(),
        "deadline_empty": make_deadline_empty(),
        "schedule_day": make_schedule_day(),
        "schedule_week": make_schedule_week(),
        "schedule_month": make_schedule_month(),
        "schedule_empty": make_schedule_empty(),
        "lowfi_home": make_lowfi("lowfi_home", "StudyFlow", ["Today overview", "Next class", "Urgent deadline"]),
        "lowfi_deadline": make_lowfi("lowfi_deadline", "Deadline", ["Search + filter", "Overdue section", "Upcoming card"]),
        "lowfi_review": make_lowfi("lowfi_review", "Review Plan", ["Goal card", "Study session", "Progress"]),
    }
    return paths


def add_bg(slide, dark=False):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY if dark else SOFT


def tb(slide, x, y, w, h, value, size=18, color=INK, bold=False,
       align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = value
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def bullet_box(slide, x, y, w, h, items, size=14, color=MUTED, bullet=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}" if bullet else item
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)
    return box


def ppt_shape(slide, x, y, w, h, fill=CARD, line=BORDER, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(0.75)
    return shp


def ppt_pill(slide, x, y, w, h, value, fill, color, size=10):
    shp = ppt_shape(slide, x, y, w, h, fill, fill)
    tf = shp.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return shp


def header(slide, index, title, subtitle="", dark=False):
    color = RGBColor(255, 255, 255) if dark else INK
    sub = RGBColor(192, 203, 220) if dark else MUTED
    tb(slide, Inches(0.65), Inches(0.38), Inches(1.0), Inches(0.24),
       f"{index:02d}", 10, BLUE, True)
    tb(slide, Inches(0.65), Inches(0.68), Inches(8.8), Inches(0.52),
       title, 27, color, True)
    if subtitle:
        tb(slide, Inches(0.65), Inches(1.25), Inches(9.8), Inches(0.32),
           subtitle, 12, sub)
    tb(slide, Inches(11.5), Inches(7.05), Inches(1.1), Inches(0.2),
       "StudyFlow", 8, sub, True, PP_ALIGN.RIGHT)


def image_card(slide, path, x, y, w, h, label=None):
    ppt_shape(slide, x, y, w, h, CARD, BORDER)
    slide.shapes.add_picture(str(path), x + Inches(0.12), y + Inches(0.12),
                             width=w - Inches(0.24), height=h - Inches(0.24))
    if label:
        ppt_pill(slide, x + Inches(0.18), y + h - Inches(0.42),
                 Inches(1.75), Inches(0.28), label, RGBColor(255, 255, 255), BLUE, 8)


def add_table(slide, x, y, w, h, rows, widths):
    tbl = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h).table
    for i, width in enumerate(widths):
        tbl.columns[i].width = width
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = value
            cell.margin_left = Pt(4)
            cell.margin_right = Pt(4)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r == 0 else CARD
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Arial"
            p.font.size = Pt(8.3 if r else 9)
            p.font.bold = r == 0
            p.font.color.rgb = RGBColor(255, 255, 255) if r == 0 else INK
    return tbl


def cover(prs, assets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, True)
    tb(s, Inches(0.75), Inches(0.78), Inches(4.2), Inches(0.25),
       "UX/UI CASE STUDY • MOBILE PRODUCT", 10, RGBColor(147, 197, 253), True)
    tb(s, Inches(0.75), Inches(1.24), Inches(6.1), Inches(0.9),
       "StudyFlow", 52, RGBColor(255, 255, 255), True)
    tb(s, Inches(0.8), Inches(2.26), Inches(5.95), Inches(0.82),
       "Ứng dụng hỗ trợ sinh viên quản lý lịch học, deadline và kế hoạch ôn tập một cách chủ động.",
       19, RGBColor(202, 213, 225))
    ppt_pill(s, Inches(0.8), Inches(3.34), Inches(1.35), Inches(0.38), "Schedule", RGBColor(30, 64, 175), RGBColor(191, 219, 254))
    ppt_pill(s, Inches(2.32), Inches(3.34), Inches(1.35), Inches(0.38), "Deadline", RGBColor(127, 29, 29), RGBColor(254, 202, 202))
    ppt_pill(s, Inches(3.84), Inches(3.34), Inches(1.25), Inches(0.38), "Review", RGBColor(20, 83, 45), RGBColor(187, 247, 208))
    s.shapes.add_picture(str(assets["schedule_day"]), Inches(7.35), Inches(0.72), width=Inches(2.22), height=Inches(4.82))
    s.shapes.add_picture(str(assets["deadline_list"]), Inches(9.55), Inches(1.12), width=Inches(2.22), height=Inches(4.82))
    tb(s, Inches(0.8), Inches(6.58), Inches(6.0), Inches(0.3),
       "Professional deck with Figma-based mobile visuals", 10, RGBColor(148, 163, 184), True)


def simple_cards(slide, cards, y=Inches(1.78)):
    for i, (h, b, fill, color) in enumerate(cards):
        x = Inches(0.72 + (i % 3) * 4.15)
        yy = y + Inches((i // 3) * 1.68)
        ppt_shape(slide, x, yy, Inches(3.65), Inches(1.18), fill, BORDER)
        ppt_pill(slide, x + Inches(0.2), yy + Inches(0.18), Inches(0.5), Inches(0.28), str(i + 1), color, RGBColor(255, 255, 255), 9)
        tb(slide, x + Inches(0.84), yy + Inches(0.18), Inches(2.55), Inches(0.23), h, 14, INK, True)
        bullet_box(slide, x + Inches(0.84), yy + Inches(0.5), Inches(2.55), Inches(0.44), [b], 10.5, MUTED, False)


def build():
    assets = generate_assets()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    cover(prs, assets)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 1, "Tổng quan dự án", "StudyFlow gom lịch học, deadline và kế hoạch ôn tập vào một trải nghiệm mobile thống nhất.")
    simple_cards(s, [
        ("Lịch học", "Xem ngày/tuần/tháng, thêm/sửa lớp, phòng học, giảng viên.", RGBColor(239, 246, 255), BLUE),
        ("Deadline", "Theo dõi bài tập, mức ưu tiên, quá hạn và tiến độ.", RGBColor(255, 241, 242), RED),
        ("Kế hoạch ôn tập", "Chia nhỏ mục tiêu học, nhắc ôn và theo dõi hoàn thành.", RGBColor(240, 253, 244), GREEN),
    ])
    image_card(s, assets["schedule_day"], Inches(1.1), Inches(4.1), Inches(1.8), Inches(2.6), "Schedule")
    image_card(s, assets["deadline_list"], Inches(3.25), Inches(4.1), Inches(1.8), Inches(2.6), "Deadline")
    image_card(s, assets["deadline_detail"], Inches(5.4), Inches(4.1), Inches(1.8), Inches(2.6), "Progress")
    tb(s, Inches(7.75), Inches(4.5), Inches(4.5), Inches(0.85),
       "Định vị: student command center giúp sinh viên biết rõ hôm nay cần học gì, nộp gì và ôn gì tiếp theo.",
       20, INK, True)

    for idx, heading, subtitle, items, accent in [
        (2, "Vấn đề", "Sinh viên cần học chủ động hơn nhưng thông tin bị phân tán.", [
            "Lịch học, bài tập, nhắc nhở và tài liệu nằm ở nhiều app khác nhau.",
            "Deadline đến gần nhưng không được ưu tiên rõ theo ngày/tuần.",
            "Ôn tập thường bị động, chỉ bắt đầu khi sát kiểm tra.",
            "Quá nhiều việc nhỏ khiến người dùng không biết bắt đầu từ đâu.",
        ], RED),
        (3, "Mục tiêu", "Tạo công cụ giúp sinh viên học có kế hoạch, đúng hạn và ít căng thẳng hơn.", [
            "Giảm việc quên lịch học, deadline và phiên ôn tập quan trọng.",
            "Hiển thị ưu tiên theo hôm nay, tuần này và quá hạn.",
            "Thêm/sửa/xóa nhanh trên mobile với ít bước thao tác.",
            "Tạo cảm giác tiến bộ qua trạng thái và phần trăm hoàn thành.",
            "Thiết kế dễ tiếp cận: chữ rõ, tương phản tốt, vùng chạm đủ lớn.",
        ], GREEN),
    ]:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(s)
        header(s, idx, heading, subtitle)
        bullet_box(s, Inches(0.85), Inches(1.82), Inches(6.6), Inches(4.5), items, 17, INK, True)
        image_card(s, assets["deadline_today"], Inches(8.25), Inches(1.25), Inches(2.1), Inches(4.65), "Figma visual")
        image_card(s, assets["schedule_week"], Inches(10.45), Inches(1.65), Inches(1.85), Inches(4.1), "Calendar")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 4, "Nghiên cứu người dùng", "Tổng hợp insight định tính cho nhóm sinh viên đại học.")
    add_table(s, Inches(0.65), Inches(1.55), Inches(12.05), Inches(3.55), [
        ["Phương pháp", "Mục đích", "Insight chính"],
        ["Interview 1:1", "Hiểu thói quen lập kế hoạch", "Sinh viên dùng nhiều công cụ nhưng bỏ sót cập nhật."],
        ["Survey nhanh", "Xác định tần suất quên deadline", "Deadline và lịch học là nhu cầu thường xuyên nhất."],
        ["Desk research", "So sánh app hiện có", "Ít app gom lịch học + deadline + ôn tập trong UX mobile gọn."],
        ["Usability test", "Kiểm tra flow thêm deadline/lịch", "CTA rõ và trạng thái quá hạn giúp người dùng hiểu nhanh."],
    ], [Inches(2.1), Inches(3.1), Inches(6.85)])
    ppt_shape(s, Inches(0.85), Inches(5.45), Inches(11.55), Inches(0.62), RGBColor(239, 246, 255), RGBColor(191, 219, 254))
    tb(s, Inches(1.08), Inches(5.64), Inches(11.0), Inches(0.2), "Ghi chú: Có thể thay bằng số liệu khảo sát thật của nhóm để tăng độ tin cậy.", 13, BLUE, True)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 5, "Phân tích đối thủ", "Hơn 20 app được nhóm theo loại cạnh tranh.")
    add_table(s, Inches(0.45), Inches(1.36), Inches(12.45), Inches(4.7), [
        ["Nhóm", "Ứng dụng", "Điểm mạnh", "Khoảng trống cho StudyFlow"],
        ["Student planner", "MyStudyLife, myHomework, Power Planner, iStudiez Pro, School Planner", "Theo dõi lớp, bài tập, nhắc nhở", "Chưa nhấn mạnh kế hoạch ôn tập cá nhân"],
        ["Deadline tracker", "Assignment Planner, Egenda, Class Timetable, Studious, Homework Planner", "Tập trung bài tập/hạn nộp", "Ít kết nối với lịch học và tiến độ ôn thi"],
        ["Productivity", "Todoist, TickTick, Microsoft To Do, Any.do, Trello", "Task mạnh, nhắc nhở, label", "Không sinh viên-hóa môn học/lớp học"],
        ["Calendar/planning", "Google Calendar, Notion, Structured, FlowSavvy, Apple Calendar", "Lịch, time blocking, template", "Cần tự cấu hình, dễ quá tải"],
        ["Study tools/LMS", "Quizlet, Anki, Forest, Focus To-Do, Google Classroom, Canvas Student", "Ôn tập, tập trung, LMS", "Không gom toàn bộ journey học tập cá nhân"],
    ], [Inches(1.45), Inches(3.55), Inches(2.85), Inches(4.6)])
    tb(s, Inches(0.75), Inches(6.28), Inches(11.8), Inches(0.28), "Cơ hội: StudyFlow tập trung vào hành trình học tập của sinh viên, không chỉ là to-do list.", 17, INK, True, PP_ALIGN.CENTER)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 6, "Pain Points", "Những nỗi đau chính được ưu tiên giải quyết trong phiên bản đầu.")
    simple_cards(s, [
        ("Quên deadline", "Không thấy mức độ khẩn cấp trước khi quá muộn.", RGBColor(255, 241, 242), RED),
        ("Không biết ưu tiên", "Không rõ hôm nay nên học gì trước.", RGBColor(255, 246, 221), ORANGE),
        ("Khó cập nhật lịch", "Lịch học thay đổi nhưng thao tác sửa mất thời gian.", RGBColor(239, 246, 255), BLUE),
        ("Bài lớn mơ hồ", "Không chia nhỏ tiến độ thành các bước làm được.", RGBColor(240, 253, 244), GREEN),
        ("Ôn tập thiếu kế hoạch", "Thiếu nhắc nhở và mục tiêu theo ngày.", RGBColor(236, 254, 255), CYAN),
        ("Giao diện quá tải", "Nhiều chữ làm người dùng bỏ qua thông tin chính.", RGBColor(253, 242, 248), PINK),
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 7, "Persona", "Đại diện nhóm sinh viên cần quản lý nhiều môn và nhiều hạn nộp.")
    ppt_shape(s, Inches(0.75), Inches(1.5), Inches(3.2), Inches(4.85), RGBColor(239, 246, 255), RGBColor(191, 219, 254))
    tb(s, Inches(1.02), Inches(1.9), Inches(2.6), Inches(0.4), "Minh Anh", 25, INK, True, PP_ALIGN.CENTER)
    tb(s, Inches(1.02), Inches(2.42), Inches(2.6), Inches(0.25), "Sinh viên năm 2 ngành CNTT", 12, MUTED, False, PP_ALIGN.CENTER)
    bullet_box(s, Inches(1.05), Inches(3.05), Inches(2.5), Inches(1.4), ["Nộp bài đúng hạn", "Biết hôm nay cần học gì", "Không ôn thi sát giờ"], 12, INK)
    tb(s, Inches(4.45), Inches(1.7), Inches(3.3), Inches(0.28), "Hành vi hiện tại", 17, INK, True)
    bullet_box(s, Inches(4.45), Inches(2.12), Inches(3.65), Inches(1.65), ["Ghi deadline trong chat nhóm hoặc Google Calendar.", "Dùng to-do list nhưng ít cập nhật tiến độ.", "Thường học theo cảm giác thay vì kế hoạch."], 13, MUTED)
    tb(s, Inches(8.4), Inches(1.7), Inches(3.3), Inches(0.28), "Nhu cầu", 17, INK, True)
    bullet_box(s, Inches(8.4), Inches(2.12), Inches(3.75), Inches(1.65), ["Một màn hình biết việc nào gấp.", "Flow thêm deadline nhanh.", "Nhắc nhở nhẹ, không gây áp lực.", "Theo dõi tiến độ rõ bằng phần trăm."], 13, MUTED)
    ppt_shape(s, Inches(4.45), Inches(4.68), Inches(7.7), Inches(1.0), RGBColor(255, 246, 221), RGBColor(253, 230, 138))
    tb(s, Inches(4.8), Inches(4.96), Inches(7.0), Inches(0.32), "“Mình không thiếu app, mình thiếu một chỗ cho biết nên làm gì tiếp theo.”", 16, INK, True, PP_ALIGN.CENTER)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 8, "User Journey", "Từ mở app đến hoàn thành một nhiệm vụ học tập.")
    steps = [("Mở app", "Biết hôm nay có gì"), ("Xem lịch", "Kiểm tra lớp và giờ trống"), ("Chọn deadline", "Biết việc gấp"), ("Lập ôn tập", "Chia nhỏ bài học"), ("Hoàn thành", "Cập nhật tiến độ")]
    for i, (h, b) in enumerate(steps):
        x = Inches(0.62 + i * 2.5)
        ppt_shape(s, x, Inches(1.95), Inches(2.08), Inches(2.6), CARD, BORDER)
        ppt_pill(s, x + Inches(0.22), Inches(2.18), Inches(0.42), Inches(0.32), str(i + 1), BLUE, RGBColor(255, 255, 255))
        tb(s, x + Inches(0.22), Inches(2.75), Inches(1.65), Inches(0.25), h, 13, INK, True)
        bullet_box(s, x + Inches(0.22), Inches(3.12), Inches(1.65), Inches(0.55), [b], 10.5, MUTED, False)
    tb(s, Inches(0.9), Inches(5.35), Inches(11.3), Inches(0.36), "Design implication: home screen phải trả lời nhanh “hôm nay cần làm gì?”", 18, INK, True, PP_ALIGN.CENTER)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 9, "Ideation", "Ý tưởng được tạo từ pain points và chuyển thành feature ưu tiên.")
    bullet_box(s, Inches(0.85), Inches(1.75), Inches(5.8), Inches(2.4), [
        "HMW giúp sinh viên thấy deadline quan trọng trước khi quá muộn?",
        "HMW kết nối lịch học với kế hoạch làm bài và ôn tập?",
        "HMW biến tiến độ học thành tín hiệu trực quan, dễ hiểu?",
    ], 15, INK)
    add_table(s, Inches(6.95), Inches(1.62), Inches(5.4), Inches(3.25), [
        ["Feature", "Giá trị"],
        ["Calendar views", "Nhìn lịch theo ngày/tuần/tháng"],
        ["Deadline badges", "Hôm nay, 2 ngày, quá hạn rõ ràng"],
        ["Progress bar", "Biết mức độ hoàn thành"],
        ["Review plan", "Ôn tập theo mục tiêu"],
    ], [Inches(2.0), Inches(3.4)])
    image_card(s, assets["deadline_today"], Inches(1.2), Inches(4.35), Inches(1.45), Inches(2.1), "Today")
    image_card(s, assets["deadline_overdue"], Inches(2.9), Inches(4.35), Inches(1.45), Inches(2.1), "Risk")
    image_card(s, assets["schedule_week"], Inches(4.6), Inches(4.35), Inches(1.45), Inches(2.1), "Week")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 10, "Low-Fi Design", "Wireframe tập trung vào cấu trúc nội dung và flow chính trước khi xử lý visual.")
    image_card(s, assets["lowfi_home"], Inches(1.0), Inches(1.45), Inches(2.2), Inches(4.75), "Home")
    image_card(s, assets["lowfi_deadline"], Inches(5.55), Inches(1.45), Inches(2.2), Inches(4.75), "Deadline")
    image_card(s, assets["lowfi_review"], Inches(10.1), Inches(1.45), Inches(2.2), Inches(4.75), "Review")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 11, "Low-Fi Prototype", "Prototype kiểm tra luồng tác vụ cốt lõi trước khi đầu tư high-fi.")
    flow = [("Home", 0.8, 2.05), ("Lịch học", 3.05, 1.45), ("Thêm lịch", 5.45, 1.45), ("Deadline", 3.05, 2.85), ("Chi tiết", 5.45, 2.85), ("Ôn tập", 3.05, 4.25), ("Kế hoạch", 5.45, 4.25), ("Hoàn thành", 8.05, 2.85)]
    for label, x, y in flow:
        ppt_shape(s, Inches(x), Inches(y), Inches(1.65), Inches(0.62), CARD, BORDER)
        tb(s, Inches(x), Inches(y + 0.2), Inches(1.65), Inches(0.18), label, 10.5, INK, True, PP_ALIGN.CENTER)
    for x1, y1, x2, y2 in [(2.45, 2.36, 3.05, 1.76), (2.45, 2.36, 3.05, 3.16), (2.45, 2.36, 3.05, 4.56), (4.7, 1.76, 5.45, 1.76), (4.7, 3.16, 5.45, 3.16), (4.7, 4.56, 5.45, 4.56), (7.1, 3.16, 8.05, 3.16)]:
        line = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = BLUE
        line.line.width = Pt(1.5)
    image_card(s, assets["lowfi_deadline"], Inches(10.05), Inches(1.75), Inches(1.85), Inches(4.0), "Wireframe")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 12, "Usability Testing", "Kiểm tra khả năng hiểu flow và thao tác nhanh trên mobile.")
    add_table(s, Inches(0.7), Inches(1.55), Inches(11.95), Inches(3.65), [
        ["Task", "Tiêu chí thành công", "Kết quả/Insight"],
        ["Thêm deadline mới", "Hoàn thành dưới 60 giây", "Cần label rõ cho ngày/giờ đến hạn"],
        ["Tìm deadline hôm nay", "Không quá 2 lần tap", "Tab Hôm nay giúp giảm thời gian tìm kiếm"],
        ["Sửa tiến độ", "Người dùng hiểu % hoàn thành", "Progress bar nên đi kèm số %"],
        ["Xem quá hạn", "Nhận biết ngay rủi ro", "Màu đỏ và badge Quá hạn hoạt động tốt"],
    ], [Inches(2.7), Inches(3.15), Inches(6.1)])
    image_card(s, assets["deadline_add"], Inches(0.95), Inches(5.45), Inches(1.1), Inches(1.35), "Add")
    image_card(s, assets["deadline_today"], Inches(2.28), Inches(5.45), Inches(1.1), Inches(1.35), "Today")
    image_card(s, assets["deadline_overdue"], Inches(3.61), Inches(5.45), Inches(1.1), Inches(1.35), "Overdue")
    tb(s, Inches(5.0), Inches(5.76), Inches(6.7), Inches(0.32), "Khuyến nghị: trạng thái rõ, CTA nổi bật, giảm nhập liệu bằng dropdown/date picker.", 14, INK, True)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 13, "Cải tiến", "Thay đổi sau vòng test để tăng rõ ràng và tốc độ thao tác.")
    simple_cards(s, [
        ("Trước", "Deadline chỉ là danh sách dài, khó biết việc nào gấp.", RGBColor(255, 241, 242), RED),
        ("Sau", "Tách Quá hạn, Sắp tới, Hôm nay, Tuần; thêm badge ngày còn lại.", RGBColor(240, 253, 244), GREEN),
        ("Trước", "Form nhập nhiều text tự do.", RGBColor(255, 241, 242), RED),
        ("Sau", "Dùng dropdown, date picker, time picker, chip ưu tiên.", RGBColor(240, 253, 244), GREEN),
    ])
    image_card(s, assets["deadline_list"], Inches(8.45), Inches(3.15), Inches(1.55), Inches(3.2), "After")
    image_card(s, assets["deadline_add"], Inches(10.25), Inches(3.15), Inches(1.55), Inches(3.2), "Form")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 14, "High-Fi Design", "Visual direction dựa trên Figma: mobile-first, card rõ, màu trạng thái dễ nhận biết.")
    image_card(s, assets["schedule_day"], Inches(0.7), Inches(1.35), Inches(1.7), Inches(4.85), "Schedule day")
    image_card(s, assets["schedule_week"], Inches(2.68), Inches(1.35), Inches(1.7), Inches(4.85), "Week")
    image_card(s, assets["deadline_list"], Inches(4.66), Inches(1.35), Inches(1.7), Inches(4.85), "Deadline")
    image_card(s, assets["deadline_detail"], Inches(6.64), Inches(1.35), Inches(1.7), Inches(4.85), "Detail")
    tb(s, Inches(9.0), Inches(1.7), Inches(3.2), Inches(0.35), "Design System", 18, INK, True)
    bullet_box(s, Inches(9.0), Inches(2.15), Inches(3.4), Inches(2.2), ["8px radius cho card/button", "Primary blue #3478F6", "Red cho quá hạn", "Green cho hoàn thành", "Typography rõ trên mobile"], 13, MUTED)
    ppt_pill(s, Inches(9.0), Inches(4.8), Inches(0.92), Inches(0.34), "#3478F6", RGBColor(239, 246, 255), BLUE, 8.5)
    ppt_pill(s, Inches(10.05), Inches(4.8), Inches(0.92), Inches(0.34), "#EF4444", RGBColor(255, 241, 242), RED, 8.5)
    ppt_pill(s, Inches(11.1), Inches(4.8), Inches(0.92), Inches(0.34), "#22C55E", RGBColor(240, 253, 244), GREEN, 8.5)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 15, "High-Fi Prototype", "Prototype hoàn thiện mô phỏng đầy đủ các trạng thái người dùng sẽ gặp.")
    prototype_imgs = [("List", assets["deadline_list"]), ("Today", assets["deadline_today"]), ("Detail", assets["deadline_detail"]), ("Add", assets["deadline_add"]), ("Edit", assets["deadline_edit"]), ("Empty", assets["deadline_empty"])]
    for i, (label, path) in enumerate(prototype_imgs):
        x = Inches(0.72 + i * 2.05)
        image_card(s, path, x, Inches(1.45), Inches(1.55), Inches(4.25), label)
        if i < len(prototype_imgs) - 1:
            line = s.shapes.add_connector(1, x + Inches(1.6), Inches(3.55), x + Inches(1.98), Inches(3.55))
            line.line.color.rgb = BLUE
            line.line.width = Pt(1.2)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 16, "Accessibility", "Thiết kế để sinh viên thao tác nhanh và đọc được trong nhiều bối cảnh.")
    simple_cards(s, [
        ("Contrast", "Màu chữ chính đậm, trạng thái dùng màu + text badge.", RGBColor(239, 246, 255), BLUE),
        ("Touch target", "Button, icon và chip ưu tiên khoảng 40-52px.", RGBColor(240, 253, 244), GREEN),
        ("Readable copy", "Nhãn ngắn: Hôm nay, Quá hạn, Thêm Deadline.", RGBColor(255, 246, 221), ORANGE),
        ("Navigation", "Bottom nav nhất quán cho các module StudyFlow.", RGBColor(239, 246, 255), BLUE),
        ("Error prevention", "Date/time picker, dropdown môn học, validation form.", RGBColor(240, 253, 244), GREEN),
        ("State feedback", "Badge, progress bar và empty state rõ.", RGBColor(255, 241, 242), RED),
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 17, "Demo mobile", "Kịch bản demo đề xuất cho phần trình bày.")
    bullet_box(s, Inches(0.85), Inches(1.65), Inches(5.8), Inches(4.2), [
        "Mở app và xem deadline chưa hoàn thành.",
        "Chuyển tab Hôm nay để thấy Quiz HTML/CSS.",
        "Mở chi tiết deadline và tăng tiến độ.",
        "Thêm một deadline mới với ngày, giờ, mức ưu tiên.",
        "Chuyển sang lịch học để xem ngày/tuần/tháng.",
        "Kết luận bằng trạng thái quá hạn và kế hoạch ôn tập.",
    ], 16, INK)
    image_card(s, assets["deadline_list"], Inches(7.1), Inches(1.32), Inches(1.7), Inches(4.8), "1")
    image_card(s, assets["deadline_detail"], Inches(9.0), Inches(1.32), Inches(1.7), Inches(4.8), "2")
    image_card(s, assets["deadline_add"], Inches(10.9), Inches(1.32), Inches(1.7), Inches(4.8), "3")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 18, "Code Implementation", "Prototype Flutter đã được tách thành các app/folder riêng để dễ phát triển tiếp.")
    add_table(s, Inches(0.65), Inches(1.5), Inches(12.05), Inches(4.35), [
        ["Module", "File/Folder", "Vai trò"],
        ["Schedule Calendar", "schedule_calendar_app", "Ngày/tuần/tháng, thêm/sửa lịch học, chi tiết lớp"],
        ["Assignments Deadlines", "assignments_deadlines_app", "Deadline list, today, week, overdue, add/edit/detail"],
        ["Model layer", "models/*.dart", "Định nghĩa dữ liệu ScheduleItem, DeadlineItem, priority, state"],
        ["Screen layer", "screens/*.dart", "UI, navigation, form validation, interaction"],
        ["Quality check", "flutter analyze/test/build web", "Đảm bảo project compile và test pass"],
    ], [Inches(2.0), Inches(3.0), Inches(7.05)])
    tb(s, Inches(0.85), Inches(6.1), Inches(11.4), Inches(0.34), "Hướng phát triển: gộp các module vào một StudyFlow app duy nhất, thêm lưu trữ local/backend và notification.", 15, INK, True, PP_ALIGN.CENTER)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 19, "Kết quả cuối", "Phiên bản hiện tại chứng minh được trải nghiệm cốt lõi của StudyFlow.")
    simple_cards(s, [
        ("Hoàn thành", "2 app Flutter độc lập cho Schedule và Deadline từ Figma.", RGBColor(240, 253, 244), GREEN),
        ("Tương tác", "Thêm/sửa/xóa, lọc, tìm kiếm, cập nhật tiến độ, empty state.", RGBColor(239, 246, 255), BLUE),
        ("Kiểm chứng", "Analyzer sạch, widget test pass, build web thành công.", RGBColor(240, 253, 244), GREEN),
        ("Giá trị", "Dashboard học tập rõ hơn và ít bỏ sót việc quan trọng.", RGBColor(255, 246, 221), ORANGE),
    ])
    image_card(s, assets["schedule_month"], Inches(8.2), Inches(3.35), Inches(1.35), Inches(3.15), "Calendar")
    image_card(s, assets["deadline_empty"], Inches(9.85), Inches(3.35), Inches(1.35), Inches(3.15), "Empty")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 20, "Reflection", "Bài học chính trong quá trình chuyển Figma thành prototype Flutter.")
    bullet_box(s, Inches(0.9), Inches(1.75), Inches(7.2), Inches(3.4), [
        "Thiết kế tốt cần trạng thái dữ liệu rõ: empty, overdue, today, upcoming.",
        "Sinh viên cần quyết định nhanh, nên ưu tiên hierarchy và badge trạng thái.",
        "Figma giúp định hình visual, nhưng code cần tách model/screen để dễ mở rộng.",
        "Nên validate với người dùng thật để thay giả định bằng số liệu đáng tin cậy.",
    ], 17, INK)
    ppt_shape(s, Inches(9.0), Inches(1.8), Inches(2.7), Inches(2.6), RGBColor(239, 246, 255), RGBColor(191, 219, 254))
    tb(s, Inches(9.25), Inches(2.35), Inches(2.2), Inches(0.45), "Key takeaway", 18, BLUE, True, PP_ALIGN.CENTER)
    tb(s, Inches(9.25), Inches(3.02), Inches(2.2), Inches(0.65), "Học tập chủ động bắt đầu từ việc biết rõ việc tiếp theo.", 15, INK, True, PP_ALIGN.CENTER)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    header(s, 21, "Next Steps", "Các bước tiếp theo để biến prototype thành sản phẩm StudyFlow hoàn chỉnh.")
    bullet_box(s, Inches(0.9), Inches(1.65), Inches(7.8), Inches(4.2), [
        "Gộp Schedule + Deadline + Review Plan vào một codebase Flutter.",
        "Thêm lưu trữ local bằng SQLite/Hive và đồng bộ cloud khi đăng nhập.",
        "Tích hợp notification cho lớp học, deadline và phiên ôn tập.",
        "Import lịch từ Google Calendar/LMS hoặc file syllabus.",
        "Gợi ý kế hoạch ôn tập tự động theo deadline, độ khó và thời gian trống.",
        "Chạy usability test vòng 2 với sinh viên thật và đo time-on-task.",
    ], 15.5, INK)
    image_card(s, assets["schedule_empty"], Inches(9.35), Inches(1.55), Inches(1.45), Inches(3.5), "Schedule")
    image_card(s, assets["deadline_empty"], Inches(10.95), Inches(1.55), Inches(1.45), Inches(3.5), "Deadline")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, True)
    tb(s, Inches(0.95), Inches(1.8), Inches(11.4), Inches(0.9), "Q&A", 60, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
    tb(s, Inches(1.55), Inches(2.95), Inches(10.2), Inches(0.4), "Cảm ơn thầy/cô và các bạn đã lắng nghe.", 20, RGBColor(202, 213, 225), False, PP_ALIGN.CENTER)
    ppt_pill(s, Inches(5.25), Inches(3.75), Inches(2.8), Inches(0.42), "StudyFlow", RGBColor(30, 64, 175), RGBColor(191, 219, 254), 14)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
